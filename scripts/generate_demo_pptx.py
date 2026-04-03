"""AI Demo Day — Shiproom v5 FINAL — The Coordination Tax narrative"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
AZ=RGBColor(0x00,0x78,0xD4);DK=RGBColor(0x00,0x20,0x50);WH=RGBColor(0xFF,0xFF,0xFF)
LG=RGBColor(0xF2,0xF2,0xF2);CY=RGBColor(0x50,0xE6,0xFF);GR=RGBColor(0x00,0xB7,0xC3)
OR=RGBColor(0xFF,0x8C,0x00);MG=RGBColor(0x60,0x60,0x60);BG=RGBColor(0xD2,0xD2,0xD2)
RD=RGBColor(0xD1,0x34,0x38);PU=RGBColor(0x88,0x44,0xCC);INK=RGBColor(0x1B,0x1B,0x1B)
prs=Presentation();prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5);SW=prs.slide_width
def bg(s,c):s.background.fill.solid();s.background.fill.fore_color.rgb=c
def R(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background();return sh
def T(s,l,t,w,h,tx,sz=18,c=INK,b=False,a=PP_ALIGN.LEFT,f="Segoe UI"):
    tb=s.shapes.add_textbox(l,t,w,h);tf=tb.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=tx;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b;p.font.name=f;p.alignment=a;return tb
def N(s,tx):s.notes_slide.notes_text_frame.text=tx

# === SLIDE 1: TITLE ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,DK)
R(s,Inches(0),Inches(0),SW,Inches(0.08),AZ);R(s,Inches(1.2),Inches(2.5),Inches(1.5),Inches(0.05),CY)
T(s,Inches(1.2),Inches(2.75),Inches(10),Inches(1.2),"Shiproom",sz=56,c=WH,b=True,f="Segoe UI Semibold")
T(s,Inches(1.2),Inches(4.0),Inches(10),Inches(1.2),"The invisible cost of knowing\nwhere your engineering org stands \u2014\nand how we eliminated it.",sz=22,c=CY,f="Segoe UI Light")
T(s,Inches(1.2),Inches(5.7),Inches(6),Inches(0.4),"[Speaker Name]  \u00b7  [Organization]  \u00b7  Program Management",sz=15,c=RGBColor(0xA0,0xC0,0xE0))
T(s,Inches(1.2),Inches(6.15),Inches(6),Inches(0.3),"[Event Name]  \u00b7  April 29, 2026",sz=13,c=RGBColor(0x80,0xA0,0xC0))
R(s,Inches(8.5),Inches(5.8),Inches(4.2),Inches(1),RGBColor(0x00,0x30,0x67))
T(s,Inches(8.7),Inches(5.9),Inches(3.8),Inches(0.8),"\u26a0 Replace with official title slide",sz=13,c=RGBColor(0xFF,0xCC,0x00))
N(s,"DELIVERY: Pause after title. Let them read it.\n\n\"I'm [Speaker Name] from [Organization]. I want to talk about a problem every large engineering org has \u2014 but nobody names.\"\n\n\"It's the invisible cost of simply knowing where you stand.\"\n\nMove to next slide immediately.")

# === SLIDE 2: THE UNNAMED PROBLEM ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,WH);R(s,Inches(0),Inches(0),SW,Inches(0.06),AZ)
T(s,Inches(1),Inches(1.5),Inches(11.3),Inches(2),"How much of your org's energy\ngoes into figuring out where it stands\nvs. actually moving forward?",sz=42,c=DK,b=True,f="Segoe UI Semibold",a=PP_ALIGN.CENTER)
R(s,Inches(3.5),Inches(4.0),Inches(6.3),Inches(1.6),RGBColor(0xFD,0xED,0xED))
T(s,Inches(3.5),Inches(4.1),Inches(6.3),Inches(0.8),"Research shows: 35\u201380%",sz=36,c=RD,b=True,a=PP_ALIGN.CENTER,f="Segoe UI Semibold")
T(s,Inches(3.5),Inches(4.85),Inches(6.3),Inches(0.6),"of knowledge worker time goes to coordination \u2014 not the work itself.",sz=18,c=RGBColor(0x80,0x30,0x30),a=PP_ALIGN.CENTER)
T(s,Inches(3.5),Inches(5.9),Inches(6.3),Inches(0.3),"Sources: MSR SPACE Framework \u00b7 Stray & Moe (Global SE study) \u00b7 Brooks, The Mythical Man-Month",sz=11,c=MG,a=PP_ALIGN.CENTER)
N(s,"TALKING POINTS (~40 seconds \u2014 THIS IS THE PIVOTAL SLIDE):\n\n[Pause. Let the question land.]\n\n\"Think about your own team for a second. How many hours this week were spent in meetings where the only purpose was to learn what's happening? Not to decide. Not to design. Just to know.\"\n\n\"Research puts a number on it. Engineers in global projects spend nearly 16 hours a week \u2014 40% of their time \u2014 in scheduled and unscheduled meetings. The broader research shows 35 to 80% of knowledge worker time goes to coordination. Not shipping. Coordinating.\"\n\n\"Microsoft's own SPACE framework acknowledges this as one of the biggest drags on developer productivity.\"\n\n\"Fred Brooks saw this in 1975. 30 people = 435 communication paths. Every person you add makes it harder to know what's happening.\"\n\n\"But here's what nobody talks about: this cost is invisible. It doesn't show up in any dashboard. No one tracks 'hours spent figuring out if we're on track.' It just feels like... everyone is busy. All the time.\"\n\n\"We felt it. For years. We just couldn't name it.\"\n\nKEY: Let this sink in. The audience should be nodding. They KNOW this feeling.")

# === SLIDE 3: THE VICIOUS CYCLE ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,WH);R(s,Inches(0),Inches(0),SW,Inches(0.06),AZ)
T(s,Inches(0.8),Inches(0.3),Inches(6),Inches(0.3),"THE CYCLE",sz=11,c=AZ,b=True,f="Segoe UI Semibold")
T(s,Inches(0.8),Inches(0.7),Inches(11),Inches(0.9),"Why it gets worse, never better.",sz=40,c=DK,b=True,f="Segoe UI Semibold")
nodes=[(Inches(1.5),Inches(2.5),"\U0001f624","People are too busy\nto update status"),(Inches(7.5),Inches(2.5),"\U0001f4c5","Meetings multiply\nto chase status"),(Inches(7.5),Inches(4.8),"\U0001f6ab","Less time to\ncapture anything"),(Inches(1.5),Inches(4.8),"\U0001f507","Knowledge stays\nin someone's head")]
for x,y,emoji,label in nodes:
    card=R(s,x,y,Inches(4),Inches(1.8),LG);card.line.color.rgb=BG;card.line.width=Pt(1)
    T(s,x+Inches(0.2),y+Inches(0.2),Inches(1),Inches(1.2),emoji,sz=36,c=AZ,a=PP_ALIGN.CENTER)
    T(s,x+Inches(1.2),y+Inches(0.3),Inches(2.5),Inches(1.2),label,sz=18,c=DK,b=True,f="Segoe UI Semibold")
T(s,Inches(5.7),Inches(2.9),Inches(1.5),Inches(0.6),"\u2192",sz=36,c=MG,a=PP_ALIGN.CENTER)
T(s,Inches(10),Inches(3.9),Inches(1.5),Inches(0.6),"\u2193",sz=36,c=MG,a=PP_ALIGN.CENTER)
T(s,Inches(5.7),Inches(5.2),Inches(1.5),Inches(0.6),"\u2190",sz=36,c=MG,a=PP_ALIGN.CENTER)
T(s,Inches(1.5),Inches(3.9),Inches(1.5),Inches(0.6),"\u2193",sz=36,c=MG,a=PP_ALIGN.CENTER)
T(s,Inches(1),Inches(6.5),Inches(11.3),Inches(0.6),"Every large org has this cycle. The question is whether you break it or keep living in it.",sz=19,c=MG,a=PP_ALIGN.CENTER,f="Segoe UI Italic")
N(s,"TALKING POINTS (~20 seconds \u2014 FAST, this is momentum):\n\n\"It's a cycle. People are too busy to put status in ADO \u2014 so someone schedules a meeting. More meetings means less time. Knowledge stays in someone's head. When they move teams, it's gone.\"\n\n\"Every large org has this cycle. Do you name it and break it, or keep living in it?\"\n\n\"We decided to break it.\"\n\nTRANSITION: Move quickly to journey. Energy building.")

# === SLIDE 4: THE JOURNEY ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,WH);R(s,Inches(0),Inches(0),SW,Inches(0.06),AZ)
T(s,Inches(0.8),Inches(0.3),Inches(6),Inches(0.3),"THE JOURNEY",sz=11,c=AZ,b=True,f="Segoe UI Semibold")
T(s,Inches(0.8),Inches(0.7),Inches(11),Inches(0.9),"How we broke the cycle.",sz=40,c=DK,b=True,f="Segoe UI Semibold")
stages=[("\u274c","Try & Fail","Months of\ntrial and error",RD),("\U0001f4dd","Document\n& Refine","6+ months of\ncollective learning",OR),("\u2699\ufe0f","Automate","Shiproom\nis here",AZ),("\U0001f9e0","Apply AI","Where it gets\nexciting",GR)]
cw=Inches(2.65);sx=Inches(0.5);gp=Inches(0.35)
for i,(emoji,title,sub,color) in enumerate(stages):
    x=sx+i*(cw+gp);y=Inches(2.2)
    card=R(s,x,y,cw,Inches(3.0),LG);card.line.color.rgb=BG;card.line.width=Pt(1)
    R(s,x,y,cw,Inches(0.1),color)
    T(s,x,y+Inches(0.3),cw,Inches(0.6),emoji,sz=42,c=color,a=PP_ALIGN.CENTER)
    T(s,x,y+Inches(1.0),cw,Inches(0.7),title,sz=20,c=DK,b=True,a=PP_ALIGN.CENTER,f="Segoe UI Semibold")
    T(s,x,y+Inches(1.8),cw,Inches(0.8),sub,sz=15,c=MG,a=PP_ALIGN.CENTER)
    if i==2:
        hl=R(s,x-Inches(0.04),y-Inches(0.04),cw+Inches(0.08),Inches(3.08),RGBColor(0xDE,0xEC,0xF9));hl.line.color.rgb=AZ;hl.line.width=Pt(3);s.shapes._spTree.remove(hl._element);s.shapes._spTree.insert(2,hl._element)
for i in range(3):
    ax=sx+(i+1)*(cw+gp)-gp+Inches(0.05);T(s,ax,Inches(3.4),Inches(0.3),Inches(0.5),"\u2192",sz=28,c=MG,a=PP_ALIGN.CENTER)
R(s,Inches(0.5),Inches(5.7),Inches(12.3),Inches(1.3),RGBColor(0xFE,0xF4,0xE5))
T(s,Inches(0.9),Inches(5.85),Inches(11.5),Inches(0.5),"The shortcut doesn't exist.",sz=22,c=RGBColor(0x6B,0x3A,0x00),b=True,f="Segoe UI Semibold")
T(s,Inches(0.9),Inches(6.35),Inches(11.5),Inches(0.5),"Automate what you understand. Not what you hope works.",sz=18,c=RGBColor(0x6B,0x3A,0x00))
N(s,"TALKING POINTS (~35 seconds):\n\n\"So how did we break it? Not with a tool. With a journey.\"\n\nSTAGE 1: \"We tried to solve it manually. Failed. Tried again. Failed. Eventually found what works.\"\n\nSTAGE 2: \"Documented it. Shared it. Others found gaps. Refined for six months.\"\n\nSTAGE 3: \"Only then did we build the tool. If we'd automated on day one, we'd have automated the wrong thing.\"\n\nSTAGE 4: \"Now \u2014 because we have structured data and proven process \u2014 AI becomes a real multiplier.\"\n\n\"The shortcut doesn't exist. Automate what you understand, not what you hope works.\"")

# === SLIDE 5: BEFORE ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,RGBColor(0xFD,0xED,0xED));R(s,Inches(0),Inches(0),SW,Inches(0.06),RD)
T(s,Inches(0.8),Inches(0.3),Inches(6),Inches(0.3),"BEFORE",sz=11,c=RD,b=True,f="Segoe UI Semibold")
T(s,Inches(0.8),Inches(0.7),Inches(11),Inches(0.9),"This was our Monday morning.",sz=44,c=RD,b=True,f="Segoe UI Semibold")
for i,(emoji,label,ph) in enumerate([("\U0001f4cb","ADO Queries","[ Insert: dense ADO results,\nhundreds of rows,\nno timeline ]"),("\U0001f4ca","Excel Trackers","[ Insert: conditional-formatted\nspreadsheet, stale dates ]"),("\U0001f4d1","Manual Decks","[ Insert: status PPT with\ntables copied from Excel ]")]):
    x=Inches(0.5)+i*Inches(4.2);y=Inches(2.1)
    card=R(s,x,y,Inches(3.9),Inches(4.7),WH);card.line.color.rgb=RGBColor(0xE0,0xB0,0xB0);card.line.width=Pt(1)
    T(s,x,y+Inches(0.15),Inches(3.9),Inches(0.5),f"{emoji}  {label}",sz=22,c=RD,b=True,a=PP_ALIGN.CENTER,f="Segoe UI Semibold")
    phbox=R(s,x+Inches(0.15),y+Inches(0.75),Inches(3.6),Inches(3.5),LG);phbox.line.color.rgb=BG;phbox.line.width=Pt(1)
    T(s,x+Inches(0.3),y+Inches(1.7),Inches(3.3),Inches(1.5),ph,sz=14,c=BG,a=PP_ALIGN.CENTER)
N(s,"TALKING POINTS (~25 seconds):\n\n\"This was our Monday morning. Every week.\"\n\n[Point to ADO] \"Hundreds of items across 30+ teams. If you wanted to know what's at risk, you stared at rows.\"\n\n[Point to Excel] \"So people built trackers. Stale the moment they were emailed.\"\n\n[Point to PPT] \"Every week, hours assembling a review deck. Obsolete by Friday. Rebuilt Monday.\"\n\n\"That was the coordination tax. We paid it every single week.\"\n\n\u26a0 INSERT YOUR REAL SCREENSHOTS.\n\nTRANSITION: \"So we asked \u2014 what if we could eliminate this entirely?\"\n\n\u25b6 After this slide, SWITCH TO LIVE SHIPROOM PORTAL.")

# === SLIDE 6: TRANSITION ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,DK);R(s,Inches(0),Inches(0),SW,Inches(0.06),CY)
T(s,Inches(1),Inches(1.8),Inches(11.3),Inches(1.5),"What if you could know \u2014\nwithout asking anyone?",sz=46,c=WH,b=True,f="Segoe UI Semibold",a=PP_ALIGN.CENTER)
R(s,Inches(1.5),Inches(4.0),Inches(4),Inches(1.2),RGBColor(0x50,0x20,0x20))
T(s,Inches(1.5),Inches(4.1),Inches(4),Inches(1),"Chasing status\nExcel \u00b7 Meetings \u00b7 Decks",sz=17,c=RGBColor(0xFF,0x99,0x99),a=PP_ALIGN.CENTER)
T(s,Inches(5.8),Inches(4.2),Inches(1.7),Inches(0.8),"\u2192",sz=56,c=CY,a=PP_ALIGN.CENTER)
R(s,Inches(7.8),Inches(4.0),Inches(4),Inches(1.2),RGBColor(0x00,0x40,0x70))
T(s,Inches(7.8),Inches(4.1),Inches(4),Inches(1),"Just knowing\nOne portal \u00b7 Live \u00b7 Visual",sz=17,c=CY,a=PP_ALIGN.CENTER)
T(s,Inches(1),Inches(6.2),Inches(11.3),Inches(0.5),"\u25b6  Switching to live Shiproom portal  \u25c0",sz=18,c=OR,b=True,a=PP_ALIGN.CENTER,f="Segoe UI Semibold")
N(s,"DELIVERY: Brief. 10 seconds. Then switch.\n\n\"What if you could know where every project stands \u2014 without asking anyone?\"\n\n\"Let me show you.\"\n\n\u25b6 SWITCH TO LIVE SHIPROOM PORTAL.\n\nSHOW (~2 minutes):\n\n1. PORTFOLIO TIMELINE (~40s): \"Every project. One view. Visual timelines that people actually maintain because they can see them.\"\n\n2. RISK + HYGIENE (~40s): \"What's at risk. Which teams need help. Before: very hard to find out. Now: instant.\"\n\n3. TRANSPARENCY (~40s): \"When one team hits trouble, others listen and learn. Self-serve instead of scheduling a meeting.\"\n\nAfter 2 min \u2192 switch back for vision.")

# === SLIDE 7: WHAT BECOMES POSSIBLE ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,WH);R(s,Inches(0),Inches(0),SW,Inches(0.06),AZ)
T(s,Inches(0.8),Inches(0.3),Inches(6),Inches(0.3),"WHAT BECOMES POSSIBLE",sz=11,c=AZ,b=True,f="Segoe UI Semibold")
T(s,Inches(0.8),Inches(0.7),Inches(11),Inches(0.9),"Once the foundation exists.",sz=42,c=DK,b=True,f="Segoe UI Semibold")
for i,(emoji,title,sub,color) in enumerate([("\U0001f4ca","Zero-Effort\nRetrospectives","AI reads the work.\nTells you what worked.",AZ),("\U0001f916","Agents That\nCapture Status","Standup \u2192 ADO.\nNo human typing.",GR),("\U0001f4cb","AI-Assisted\nPlanning","History informs\nfuture commitments.",PU)]):
    x=Inches(0.4)+i*Inches(4.2);y=Inches(2.2)
    card=R(s,x,y,Inches(3.9),Inches(3.2),LG);card.line.color.rgb=BG;card.line.width=Pt(1);R(s,x,y,Inches(3.9),Inches(0.1),color)
    T(s,x,y+Inches(0.3),Inches(3.9),Inches(0.6),emoji,sz=48,c=color,a=PP_ALIGN.CENTER)
    T(s,x,y+Inches(1.0),Inches(3.9),Inches(0.8),title,sz=21,c=DK,b=True,a=PP_ALIGN.CENTER,f="Segoe UI Semibold")
    T(s,x+Inches(0.3),y+Inches(2.0),Inches(3.3),Inches(0.8),sub,sz=16,c=MG,a=PP_ALIGN.CENTER)
R(s,Inches(0.4),Inches(5.8),Inches(12.5),Inches(1.3),RGBColor(0xE8,0xF4,0xFD))
T(s,Inches(0.8),Inches(6.0),Inches(11.7),Inches(0.8),"This is just the starting point. The foundation makes everything else possible.",sz=20,c=DK,a=PP_ALIGN.CENTER)
N(s,"TALKING POINTS (~60 seconds):\n\n\"Once you have the foundation \u2014 ADO as source of truth, automation ensuring quality \u2014 AI becomes practical, not aspirational.\"\n\nRETROSPECTIVES: \"We used to do retros every semester. So heavy we had to stop. Now \u2014 zero human effort. AI reads ADO, tells you what worked, what slipped, why. Spots patterns across teams.\"\n\nAGENTS: \"Imagine standup is listened to by an agent. Someone says 'my epic is delayed' \u2014 agent writes it into ADO. Nobody types status. Ever.\"\n\nPLANNING: \"Planning is the biggest time sink. With structured historical data, AI can inform and eventually drive the planning process.\"\n\nCLOSE: \"This is just the starting point. We're at Stage 3. These ideas are Stage 4. The foundation makes all of it possible.\"")

# === SLIDE 8: TAKEAWAY ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,DK);R(s,Inches(0),Inches(0),SW,Inches(0.06),CY)
T(s,Inches(1.5),Inches(1.0),Inches(10.3),Inches(2.5),"Name the problem.\nBreak the cycle.\nBuild the foundation.\nThen let AI multiply it.",sz=40,c=WH,b=True,f="Segoe UI Semibold",a=PP_ALIGN.CENTER)
R(s,Inches(2.5),Inches(4.2),Inches(8.3),Inches(1.4),RGBColor(0x00,0x30,0x67))
T(s,Inches(2.5),Inches(4.4),Inches(8.3),Inches(1),"The org that ships fastest isn't the one with the most engineers.\nIt's the one that spends the least energy figuring out where it stands.",sz=20,c=CY,a=PP_ALIGN.CENTER,f="Segoe UI Light")
N(s,"DELIVERY: Slow. Let each line land.\n\n\"Name the problem \u2014 the coordination tax, the invisible cost of knowing.\"\n\"Break the cycle \u2014 stop accepting that everyone is just busy.\"\n\"Build the foundation \u2014 source of truth, proven process, then automation.\"\n\"Then let AI multiply it.\"\n\n[Pause]\n\n\"The org that ships fastest isn't the one with the most engineers. It's the one that spends the least energy figuring out where it stands.\"\n\n\"That's the story of Shiproom. Thank you.\"")

# === SLIDE 9: CLOSING ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,DK);R(s,Inches(0),Inches(0),SW,Inches(0.06),CY)
T(s,Inches(1),Inches(1.8),Inches(11),Inches(0.8),"Thank You",sz=48,c=WH,b=True,f="Segoe UI Semibold",a=PP_ALIGN.CENTER)
T(s,Inches(1),Inches(2.9),Inches(11),Inches(0.6),"Come see the live demo at the booth.",sz=22,c=CY,a=PP_ALIGN.CENTER,f="Segoe UI Light")
T(s,Inches(1),Inches(4.0),Inches(11),Inches(0.4),"[Speaker Name]  \u00b7  [Organization]  \u00b7  Program Management",sz=17,c=RGBColor(0xA0,0xC0,0xE0),a=PP_ALIGN.CENTER)
T(s,Inches(1),Inches(4.5),Inches(11),Inches(0.4),"[email]  \u00b7  Booth [TBD]",sz=15,c=RGBColor(0x80,0xA0,0xC0),a=PP_ALIGN.CENTER)
R(s,Inches(5.5),Inches(5.3),Inches(1.6),Inches(1.6),WH)
T(s,Inches(5.5),Inches(5.65),Inches(1.6),Inches(0.6),"[ QR ]",sz=22,c=MG,a=PP_ALIGN.CENTER)
T(s,Inches(3),Inches(5.6),Inches(2.3),Inches(0.6),"Scan to\nlearn more \u2192",sz=15,c=RGBColor(0xA0,0xC0,0xE0),a=PP_ALIGN.RIGHT)
R(s,Inches(8.5),Inches(5.8),Inches(4.2),Inches(1),RGBColor(0x00,0x30,0x67))
T(s,Inches(8.7),Inches(5.9),Inches(3.8),Inches(0.8),"\u26a0 Replace with official closing slide",sz=13,c=RGBColor(0xFF,0xCC,0x00))
N(s,"\"Come visit the booth. Happy to talk. Thank you.\"\n\nBOOTH DEPTH:\n\nCOORDINATION TAX: \"Research shows 35-80% of knowledge worker time goes to coordination. We broke that cycle.\"\n\nACCOUNTABILITY: \"When one team hits trouble, others listen and learn. Transparency through empathy, not blame.\"\n\nORG MEMORY: \"Some people carry decades of knowledge. When they move, it's gone. ADO preserves it.\"\n\nCUSTOMER ENGAGEMENT: \"Same problem with customer conversations. Light structure in existing tools \u2014 not CRM \u2014 lets AI generate reports.\"\n\nFOR SENIOR LEADERS: \"This isn't about Shiproom. It's about whether your org has the foundation to benefit from AI at all.\"")

# === SAVE ===
out=os.path.join(os.path.expanduser("~"),"OneDrive","Artha","docs","AI-Demo-Day-Shiproom-v5.pptx")
prs.save(out)
print(f"\u2705 Saved: {out}")
print(f"   Slides: {len(prs.slides)} | All with speaker notes")
print(f"\n   NARRATIVE ARC:")
print(f"     1. Title \u2014 'The invisible cost of knowing'")
print(f"     2. The Unnamed Problem \u2014 '35-80% goes to coordination'")
print(f"     3. The Vicious Cycle \u2014 visual cycle diagram")
print(f"     4. The Journey \u2014 'How we broke the cycle'")
print(f"     5. Before \u2014 ADO/Excel/PPT (your screenshots)")
print(f"     6. Transition \u2014 'What if you could know without asking?'")
print(f"     7. What Becomes Possible \u2014 retros, agents, planning")
print(f"     8. Takeaway \u2014 'Name it. Break it. Build it. Multiply it.'")
print(f"     9. Closing + QR code")
