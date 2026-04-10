#!/usr/bin/env python3
"""scripts/local_docs_kb_sync.py — Local OneDrive / SharePoint sync folder → KB sync.

Scans locally synced files (no Graph API / Entra token required) and ingests them
into the knowledge graph using the same pipeline as sharepoint_kb_sync.py:

    scan files  →  extract text  →  add_episode()  →  DocumentExtractor
                →  upsert_entity()  →  write_markdown_stub()

Covers:
  - Personal OneDrive:  C:\\Users\\<user>\\OneDrive\\
  - Synced SharePoint sites: any additional roots you configure (see --add-root)
  - Excludes the Artha workspace itself to avoid a meta-loop

Why use this instead of sharepoint_kb_sync.py?
  - No Azure App Registration / Entra config needed
  - No token refresh dance
  - Works offline
  - "Shared With Me" links that you synced locally are included

Limitation: Files shared *only* via a link (never synced to disk) are not
reachable here. Use sharepoint_kb_sync.py + Graph API for those.

Usage:
    python scripts/local_docs_kb_sync.py [--dry-run] [--verbose]
    python scripts/local_docs_kb_sync.py --lookback-days 90
    python scripts/local_docs_kb_sync.py --add-root "C:/Users/you/Contoso"
    python scripts/local_docs_kb_sync.py --status
    python scripts/local_docs_kb_sync.py --full-sync   # re-ingest all files

Supported formats: .docx  .pptx  .txt  .md  (others: skipped gracefully)

Ref: specs/kb-graph.md §10.7
"""
from __future__ import annotations

import argparse
import hashlib
import json
import logging
import os
import re
import sys
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Iterator

# ---------------------------------------------------------------------------
# Path setup
# ---------------------------------------------------------------------------
_SCRIPTS = Path(__file__).resolve().parent
_ARTHA   = _SCRIPTS.parent
if str(_SCRIPTS) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS))

from lib.knowledge_graph import get_kb          # noqa: E402
from lib.document_extractor import (            # noqa: E402
    DocumentExtractor,
    write_markdown_stub,
)

_log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
_STATE_FILE    = _ARTHA / "state" / "connectors" / "local_docs_state.json"
_NOTES_DIR     = _ARTHA / "state" / "connectors" / "local_docs_notes"
_ALLOWED_EXTS  = frozenset({".docx", ".pptx", ".txt", ".md"})
_DEFAULT_DAYS  = 90
_MAX_FILE_MB   = 20  # skip files larger than this

# Domain inference (mirrors sharepoint_kb_sync.py)
_DOMAIN_KEYWORDS: list[tuple[str, str]] = [
    ("finance",     "finance"),
    ("budget",      "finance"),
    ("legal",       "legal"),
    ("contract",    "legal"),
    ("hr",          "hr"),
    ("people",      "people"),
    ("engineering", "engineering"),
    ("dev",         "engineering"),
    ("arch",        "engineering"),
    ("design",      "engineering"),
    ("infra",       "infrastructure"),
    ("deploy",      "deployment"),
    ("ops",         "operations"),
    ("platform",    "platform"),
    ("product",     "product"),
    ("roadmap",     "product"),
    ("strategy",    "strategy"),
    ("executive",   "strategy"),
    ("meeting",     "operations"),
    ("notes",       "operations"),
    ("wiki",        "knowledge"),
    ("knowledge",   "knowledge"),
    ("learn",       "learning"),
    ("training",    "learning"),
    ("work",        "work"),
]

# ---------------------------------------------------------------------------
# Default scan roots — personal OneDrive on this machine
# ---------------------------------------------------------------------------

def _default_scan_roots() -> list[Path]:
    """Return default scan roots based on the current platform."""
    roots: list[Path] = []
    if os.name == "nt":
        onedrive = Path(os.environ.get("OneDrive", Path.home() / "OneDrive"))
    else:
        onedrive = Path.home() / "OneDrive"
    if onedrive.is_dir():
        roots.append(onedrive)
    return roots


# ---------------------------------------------------------------------------
# Text extraction (no Graph API needed — reads local files directly)
# ---------------------------------------------------------------------------

def _extract_text_from_path(path: Path) -> str:
    """Extract plain text from a local file. Returns '' on failure."""
    suffix = path.suffix.lower()
    try:
        if suffix in (".txt", ".md"):
            return path.read_text(encoding="utf-8", errors="replace")

        content = path.read_bytes()

        if suffix == ".docx":
            return _extract_docx(content, path.name)

        if suffix == ".pptx":
            return _extract_pptx(content, path.name)

    except PermissionError:
        _log.debug("Permission denied: %s", path)
    except OSError as exc:
        _log.debug("Cannot read %s: %s", path, exc)
    except Exception as exc:
        _log.warning("Failed to extract text from %s: %s", path.name, exc)
    return ""


_OLE2_MAGIC = b"\xd0\xcf\x11\xe0"  # Word 97-2003 binary format (.doc)


def _extract_docx(content: bytes, filename: str) -> str:
    """Extract text from .docx bytes via python-docx."""
    # Detect legacy binary .doc files (OLE2) with .docx extension — skip silently
    if content[:4] == _OLE2_MAGIC:
        _log.debug("Skipping legacy OLE2 binary format (not a real .docx): %s", filename)
        return ""
    try:
        from docx import Document  # type: ignore
    except ImportError:
        _log.warning(
            "python-docx not installed — %s skipped (pip install python-docx)", filename
        )
        return ""
    import io
    try:
        doc = Document(io.BytesIO(content))
    except Exception as exc:
        _log.debug("Cannot parse %s (bad or unsupported format): %s", filename, exc)
        return ""
    lines: list[str] = []
    for para in doc.paragraphs:
        if para.style and para.style.name.startswith("Heading"):
            level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 2
            lines.append(f"{'#' * level} {para.text}")
        elif para.text.strip():
            lines.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines)


def _extract_pptx(content: bytes, filename: str) -> str:
    """Extract text from .pptx bytes via python-pptx."""
    # Detect legacy binary .ppt files (OLE2) with .pptx extension — skip silently
    if content[:4] == _OLE2_MAGIC:
        _log.debug("Skipping legacy OLE2 binary format (not a real .pptx): %s", filename)
        return ""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        import traceback as _tb
        _log.debug("python-pptx ImportError for %s:\n%s", filename, _tb.format_exc())
        _log.warning(
            "python-pptx not installed — %s skipped (pip install python-pptx)", filename
        )
        return ""
    import io
    try:
        prs = Presentation(io.BytesIO(content))
    except Exception as exc:
        _log.debug("Cannot parse %s (bad or unsupported format): %s", filename, exc)
        return ""
    lines: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_lines.append(text)
        if slide_lines:
            lines.append(f"## Slide {slide_num}")
            lines.extend(slide_lines)
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# State management (lightweight local etag: path + mtime + size)
# ---------------------------------------------------------------------------

def _load_state() -> dict:
    if not _STATE_FILE.exists():
        return {"ingested": {}, "last_run_at": None}
    try:
        return json.loads(_STATE_FILE.read_text(encoding="utf-8"))
    except Exception:
        return {"ingested": {}, "last_run_at": None}


def _save_state(state: dict) -> None:
    _STATE_FILE.parent.mkdir(parents=True, exist_ok=True)
    tmp = _STATE_FILE.with_suffix(".json.tmp")
    tmp.write_text(json.dumps(state, indent=2), encoding="utf-8")
    tmp.replace(_STATE_FILE)


def _file_etag(path: Path) -> str:
    stat = path.stat()
    return f"{stat.st_size}:{stat.st_mtime_ns}"


# ---------------------------------------------------------------------------
# File discovery
# ---------------------------------------------------------------------------

def _iter_files(
    roots: list[Path],
    since: datetime,
    exclude: list[Path],
) -> Iterator[Path]:
    """Yield eligible files from roots, respecting lookback window and exclusions."""
    exclude_resolved = {p.resolve() for p in exclude}

    for root in roots:
        if not root.is_dir():
            _log.warning("Scan root not found: %s", root)
            continue
        for path in root.rglob("*"):
            if not path.is_file():
                continue
            # Skip excluded directories (e.g. the Artha workspace)
            if any(path.resolve().is_relative_to(ex) for ex in exclude_resolved):
                continue
            if path.suffix.lower() not in _ALLOWED_EXTS:
                continue
            # Skip cloud-only placeholder files (not downloaded locally)
            if _is_cloud_placeholder(path):
                _log.debug("Skipping cloud placeholder: %s", path.name)
                continue
            # Size guard
            try:
                size_mb = path.stat().st_size / (1024 * 1024)
            except OSError:
                continue
            if size_mb > _MAX_FILE_MB:
                _log.debug("Skipping large file (%.1f MB): %s", size_mb, path.name)
                continue
            # Modification time window
            try:
                mtime = datetime.fromtimestamp(path.stat().st_mtime, tz=timezone.utc)
            except OSError:
                continue
            if mtime < since:
                continue
            yield path


# ---------------------------------------------------------------------------
# Domain inference
# ---------------------------------------------------------------------------

def _infer_domain(path: Path) -> str:
    # Use only the last 4 path components and require word-boundary match
    # to avoid false-positives like 'hr' matching inside 'vemishra' or path separators.
    parts = list(path.parts[-4:])
    haystack = " ".join(p.lower() for p in parts)
    for kw, domain in _DOMAIN_KEYWORDS:
        # Match keyword as a whole word only (avoid 'hr' matching inside 'vemishra')
        if re.search(r"(?<![a-z0-9])" + re.escape(kw) + r"(?![a-z0-9])", haystack):
            return domain
    return "knowledge"


def _safe_stem(filename: str) -> str:
    stem = Path(filename).stem
    stem = re.sub(r"[^A-Za-z0-9_\-]", "_", stem)
    return stem[:60].strip("_") or "local_doc"


# ---------------------------------------------------------------------------
# Cloud-placeholder detection (Windows OneDrive)
# ---------------------------------------------------------------------------

# FILE_ATTRIBUTE_RECALL_ON_DATA_ACCESS (0x400000) indicates a cloud-only file
# that will trigger a download when read. Skip these to avoid slow/failed reads.
_WIN_RECALL_FLAG = 0x400000


def _is_cloud_placeholder(path: Path) -> bool:
    """Return True if path is a Windows OneDrive cloud-only placeholder (not synced)."""
    if os.name != "nt":
        return False
    try:
        attrs = path.stat().st_file_attributes  # type: ignore[attr-defined]
        return bool(attrs & _WIN_RECALL_FLAG)
    except (AttributeError, OSError):
        return False


# ---------------------------------------------------------------------------
# PII gate
# ---------------------------------------------------------------------------

def _pii_clean(text: str, filename: str) -> str | None:
    """Return text if PII-safe, None if blocked. Warns but does not raise."""
    try:
        from pii_guard import scan  # type: ignore
        pii_found, _ = scan(text)
        if pii_found:
            _log.warning("PII detected in '%s' — skipping ingestion", filename)
            return None
        return text
    except ImportError:
        return text  # pii_guard absent — passthrough (acceptable for local personal files)
    except Exception as exc:
        _log.warning("pii_guard error for '%s': %s — skipping to be safe", filename, exc)
        return None


# ---------------------------------------------------------------------------
# Main sync routine
# ---------------------------------------------------------------------------

def sync(
    *,
    roots: list[Path],
    lookback_days: int = _DEFAULT_DAYS,
    dry_run: bool = False,
    verbose: bool = False,
    full_sync: bool = False,
) -> dict:
    """Scan local files and ingest new/changed docs into the KB. Returns stats."""
    stats: dict = {
        "files_scanned":    0,
        "files_skipped":    0,
        "docs_ingested":    0,
        "entities_added":   0,
        "episodes_created": 0,
        "stubs_written":    0,
        "errors":           [],
    }

    since  = datetime.now(timezone.utc) - timedelta(days=lookback_days)
    state  = _load_state() if not full_sync else {"ingested": {}, "last_run_at": None}
    exclude = [_ARTHA]   # never ingest Artha's own workspace
    _NOTES_DIR.mkdir(parents=True, exist_ok=True)

    kg = None if dry_run else get_kb(_ARTHA)

    for path in _iter_files(roots, since=since, exclude=exclude):
        stats["files_scanned"] += 1
        etag = _file_etag(path)
        key  = str(path.resolve())

        # Skip unchanged files (unless full_sync)
        if not full_sync and state["ingested"].get(key) == etag:
            stats["files_skipped"] += 1
            if verbose:
                _log.info("  [skip unchanged] %s", path.name)
            continue

        content_text = _extract_text_from_path(path)
        if not content_text.strip():
            _log.debug("Empty extraction: %s", path.name)
            state["ingested"][key] = etag
            continue

        clean_text = _pii_clean(content_text, path.name)
        if clean_text is None:
            stats["files_skipped"] += 1
            continue

        domain       = _infer_domain(path)
        content_hash = hashlib.sha256(clean_text.encode()).hexdigest()
        episode_key  = f"local_docs:{content_hash[:16]}"
        source_ref   = path.as_uri()
        name         = path.name

        if verbose:
            _log.info(
                "  [%s] %s (%d chars, domain=%s)",
                path.suffix[1:], name, len(clean_text), domain,
            )

        if dry_run:
            try:
                extractor = DocumentExtractor.from_text(
                    clean_text, source_ref=source_ref, domain=domain,
                    episode_key=episode_key,
                )
                extractor.run()
                stats["docs_ingested"]  += 1
                stats["entities_added"] += len(extractor.entities)
            except Exception as exc:
                _log.warning("DRY-RUN extract error for '%s': %s", name, exc)
                stats["errors"].append(f"[dry-run] {name}: {exc}")
            continue

        try:
            # Step 1: add_episode (idempotent)
            ep_id = kg.add_episode(
                episode_key=episode_key,
                source_type="local_docs",
                raw_content=clean_text[:2000],
            )
            stats["episodes_created"] += 1

            # Step 2: extract entities
            extractor = DocumentExtractor.from_text(
                clean_text, source_ref=source_ref, domain=domain,
                episode_key=episode_key,
            )
            extractor.run()

            # Step 3: upsert entities
            for entity in extractor.entities:
                entity.setdefault("source_type", "local_docs")
                kg.upsert_entity(
                    entity,
                    source=source_ref,
                    confidence=float(entity.get("confidence", 0.65)),
                    source_episode_id=ep_id,
                )
                stats["entities_added"] += 1

            # Step 4: write markdown stub
            safe_stem = _safe_stem(name)
            hash8     = content_hash[:8]
            stub_path = _NOTES_DIR / f"{safe_stem}-{hash8}.md"
            write_markdown_stub(
                stub_path,
                title=name,
                source_ref=source_ref,
                source_type="local_docs",
                content_hash=content_hash,
                domain=domain,
                extracted_text=clean_text,
                entities=extractor.entities,
                script_name="local_docs_kb_sync.py",
            )
            stats["stubs_written"] += 1
            stats["docs_ingested"] += 1

            # Mark as ingested
            state["ingested"][key] = etag

        except Exception as exc:
            _log.error("Failed to ingest '%s': %s", name, exc)
            stats["errors"].append(f"{name}: {exc}")
            if verbose:
                import traceback
                traceback.print_exc()

    if not dry_run:
        state["last_run_at"] = datetime.now(timezone.utc).isoformat(timespec="seconds")
        _save_state(state)
        try:
            kg.close()
        except Exception:
            pass

    return stats


# ---------------------------------------------------------------------------
# Status display
# ---------------------------------------------------------------------------

def show_status() -> None:
    notes  = list(_NOTES_DIR.glob("*.md")) if _NOTES_DIR.exists() else []
    state  = _load_state()
    ingested = len(state.get("ingested", {}))
    last_run = state.get("last_run_at", "(never)")
    print(f"Local docs notes stubs: {len(notes)}")
    print(f"Notes directory:        {_NOTES_DIR}")
    print(f"State file:             {_STATE_FILE}")
    print(f"Ingested file records:  {ingested}")
    print(f"Last run:               {last_run}")

    roots = _default_scan_roots()
    print(f"\nDefault scan roots ({len(roots)}):")
    for r in roots:
        exists = "✓" if r.is_dir() else "✗ NOT FOUND"
        print(f"  {exists} {r}")
    print(f"\nArtha workspace excluded: {_ARTHA}")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        description="Artha local docs → KB sync (no Graph API needed)",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    p.add_argument(
        "--dry-run", action="store_true",
        help="Extract and parse but do not write to KB or state",
    )
    p.add_argument(
        "--full-sync", action="store_true",
        help="Re-ingest all files (ignore cached state)",
    )
    p.add_argument(
        "--status", action="store_true",
        help="Show current sync state and exit",
    )
    p.add_argument(
        "--verbose", "-v", action="store_true",
        help="Enable INFO logging",
    )
    p.add_argument(
        "--lookback-days", type=int, default=_DEFAULT_DAYS, metavar="N",
        help=f"Only scan files modified in the last N days (default: {_DEFAULT_DAYS})",
    )
    p.add_argument(
        "--add-root", dest="extra_roots", action="append", default=[], metavar="PATH",
        help="Add an extra scan root (e.g. a synced SharePoint site folder). Repeatable.",
    )
    return p


def main() -> None:
    p    = _build_parser()
    args = p.parse_args()

    level = logging.INFO if args.verbose else logging.WARNING
    logging.basicConfig(level=level, format="%(levelname)-8s %(message)s")

    if args.status:
        show_status()
        return

    roots = _default_scan_roots() + [Path(r) for r in args.extra_roots]
    if not roots:
        print("No scan roots found. OneDrive folder not detected.")
        print("Use --add-root PATH to specify a folder to scan.")
        sys.exit(1)

    print(f"Scanning {len(roots)} root(s), lookback={args.lookback_days}d ...")
    if args.dry_run:
        print("DRY RUN — no changes will be written.")

    stats = sync(
        roots=roots,
        lookback_days=args.lookback_days,
        dry_run=args.dry_run,
        verbose=args.verbose,
        full_sync=args.full_sync,
    )

    prefix = "DRY-RUN " if args.dry_run else ""
    print(f"\nLocal docs KB sync {prefix}complete:")
    for k, v in stats.items():
        if k == "errors":
            continue
        print(f"  {k}: {v}")
    if stats.get("errors"):
        print(f"\n  Errors ({len(stats['errors'])}):")
        for e in stats["errors"][:20]:
            print(f"    - {e}")
    if stats["errors"]:
        sys.exit(1)


if __name__ == "__main__":
    main()
