#!/usr/bin/env python3
# pii-guard: ignore-file — MCP server scaffolding; no personal data stored here
"""
scripts/mcp_sharepoint_server.py — Work documents MCP server for Artha.

Reads directly from the local OneDrive for Business sync folder
(C:\\Users\\<user>\\OneDrive - Microsoft\\) — no Graph API, no App
Registration, no tokens required.

Exposes three tools to Work OS agents (artha-work-msft, artha-work-enterprise):

  list_recent_documents  — list the N most recently modified work docs
  search_documents       — search by filename and/or content keywords
  read_document          — extract plain text from a local work file

Transport: stdio JSON-RPC 2.0 (MCP protocol).
"""
from __future__ import annotations

import json
import logging
import os
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

logging.basicConfig(
    level=logging.WARNING,
    format="%(levelname)s %(name)s: %(message)s",
    stream=sys.stderr,
)
_log = logging.getLogger("mcp_sharepoint_server")

# ---------------------------------------------------------------------------
# Local OneDrive root discovery
# ---------------------------------------------------------------------------

def _find_work_onedrive() -> Path | None:
    """Return the local OneDrive for Business sync root, or None."""
    user_profile = Path(os.environ.get("USERPROFILE", Path.home()))
    candidates = [
        user_profile / "OneDrive - Microsoft",
        user_profile / "OneDrive - Microsoft Corporation",
    ]
    od_env = os.environ.get("OneDriveCommercial") or os.environ.get("OneDrive")
    if od_env:
        candidates.insert(0, Path(od_env))
    for p in candidates:
        if p.exists() and p.is_dir():
            return p
    return None


_WORK_ROOT = _find_work_onedrive()

_TEXT_EXTS  = {".txt", ".md", ".csv", ".log", ".py", ".json", ".yaml", ".yml"}
_DOC_EXTS   = {".docx", ".pptx", ".xlsx"}
_OTHER_EXTS = {".pdf", ".one"}
_ALL_EXTS   = _TEXT_EXTS | _DOC_EXTS | _OTHER_EXTS
_SKIP_DIRS  = {".git", "__pycache__", ".tmp", "AppData", "~recycle"}


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def _extract_text_from_file(path: Path, max_chars: int = 10_000) -> str:
    ext = path.suffix.lower()
    if ext in _TEXT_EXTS:
        try:
            return path.read_text(encoding="utf-8", errors="replace")[:max_chars]
        except Exception:
            return ""
    if ext == ".docx":
        try:
            import docx
            doc = docx.Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())[:max_chars]
        except ImportError:
            return "[python-docx not installed — pip install python-docx]"
        except Exception as exc:
            return f"[Could not read .docx: {exc}]"
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            return "\n".join(parts)[:max_chars]
        except ImportError:
            return "[python-pptx not installed — pip install python-pptx]"
        except Exception as exc:
            return f"[Could not read .pptx: {exc}]"
    if ext == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    line = "\t".join(str(v) for v in row if v is not None)
                    if line.strip():
                        rows.append(line)
            return "\n".join(rows)[:max_chars]
        except ImportError:
            return "[openpyxl not installed — pip install openpyxl]"
        except Exception as exc:
            return f"[Could not read .xlsx: {exc}]"
    if ext == ".pdf":
        try:
            import pypdf
            reader = pypdf.PdfReader(str(path))
            return "\n".join(page.extract_text() or "" for page in reader.pages)[:max_chars]
        except ImportError:
            return "[pypdf not installed — pip install pypdf]"
        except Exception as exc:
            return f"[Could not read .pdf: {exc}]"
    return f"[Binary file — {path.suffix} — text extraction not supported]"


def _file_meta(path: Path) -> dict:
    stat = path.stat()
    rel  = str(path).replace(str(_WORK_ROOT), "").lstrip("/\\")
    return {
        "name":         path.name,
        "path":         str(path),
        "relativePath": rel,
        "ext":          path.suffix.lower(),
        "sizeBytes":    stat.st_size,
        "modified":     datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc).isoformat(),
    }


# ---------------------------------------------------------------------------
# Tool implementations
# ---------------------------------------------------------------------------

def _tool_list_recent_documents(args: dict) -> dict:
    if not _WORK_ROOT:
        return {"error": "OneDrive for Business folder not found. Is OneDrive running?"}
    limit      = min(int(args.get("limit", 20)), 100)
    ext_filter = {e.lower() for e in args.get("extensions", [])} or _ALL_EXTS
    files: list[tuple[float, Path]] = []
    try:
        for p in _WORK_ROOT.rglob("*"):
            if p.is_file() and p.suffix.lower() in ext_filter:
                if not any(skip in p.parts for skip in _SKIP_DIRS):
                    try:
                        files.append((p.stat().st_mtime, p))
                    except OSError:
                        pass
    except Exception as exc:
        return {"error": f"Could not scan OneDrive folder: {exc}"}
    files.sort(key=lambda x: x[0], reverse=True)
    results = [_file_meta(p) for _, p in files[:limit]]
    return {"files": results, "count": len(results), "root": str(_WORK_ROOT)}


def _tool_search_documents(args: dict) -> dict:
    if not _WORK_ROOT:
        return {"error": "OneDrive for Business folder not found. Is OneDrive running?"}
    query = args.get("query", "").strip().lower()
    if not query:
        return {"error": "query parameter is required"}
    limit          = min(int(args.get("limit", 15)), 50)
    content_search = bool(args.get("content_search", True))
    ext_filter     = {e.lower() for e in args.get("extensions", [])} or _ALL_EXTS
    terms = [t for t in query.split() if t]
    hits: list[dict] = []
    try:
        for path in _WORK_ROOT.rglob("*"):
            if not path.is_file():
                continue
            if path.suffix.lower() not in ext_filter:
                continue
            if any(skip in path.parts for skip in _SKIP_DIRS):
                continue
            name_lower  = path.name.lower()
            name_match  = all(t in name_lower for t in terms)
            content_match   = False
            matched_lines: list[str] = []
            if not name_match and content_search and path.suffix.lower() in _TEXT_EXTS:
                try:
                    text = path.read_text(encoding="utf-8", errors="replace")
                    if all(t in text.lower() for t in terms):
                        content_match = True
                        for line in text.splitlines():
                            if any(t in line.lower() for t in terms):
                                stripped = line.strip()
                                if stripped:
                                    matched_lines.append(stripped[:200])
                                if len(matched_lines) >= 3:
                                    break
                except OSError:
                    pass
            if name_match or content_match:
                meta = _file_meta(path)
                meta["matchType"] = "filename" if name_match else "content"
                if matched_lines:
                    meta["snippets"] = matched_lines
                hits.append(meta)
                if len(hits) >= limit:
                    break
    except Exception as exc:
        return {"error": f"Search failed: {exc}"}
    fn_hits = sorted([h for h in hits if h["matchType"] == "filename"], key=lambda h: h["modified"], reverse=True)
    ct_hits = sorted([h for h in hits if h["matchType"] == "content"], key=lambda h: h["modified"], reverse=True)
    return {"results": fn_hits + ct_hits, "count": len(hits), "query": query}


def _tool_read_document(args: dict) -> dict:
    file_path = args.get("path", "").strip()
    if not file_path:
        return {"error": "path parameter is required"}
    path = Path(file_path)
    if not path.exists():
        return {"error": f"File not found: {file_path}"}
    if not path.is_file():
        return {"error": f"Not a file: {file_path}"}
    if _WORK_ROOT:
        try:
            path.resolve().relative_to(_WORK_ROOT.resolve())
        except ValueError:
            return {"error": "Access denied: path is outside the OneDrive for Business folder"}
    max_chars = int(args.get("max_chars", 8000))
    text = _extract_text_from_file(path, max_chars=max_chars)
    return {
        "name":      path.name,
        "path":      str(path),
        "ext":       path.suffix.lower(),
        "text":      text,
        "truncated": len(text) >= max_chars,
        "chars":     len(text),
    }


# ---------------------------------------------------------------------------
# MCP protocol (stdio JSON-RPC 2.0)
# ---------------------------------------------------------------------------

_TOOLS = {
    "list_recent_documents": {
        "description": (
            "List the N most recently modified work documents from the local "
            "OneDrive for Business sync folder (C:\\Users\\...\\OneDrive - Microsoft). "
            "No authentication required. Returns paths, sizes, and timestamps."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "limit": {"type": "integer", "description": "Max files to return (default 20, max 100).", "default": 20},
                "extensions": {"type": "array", "items": {"type": "string"}, "description": "Filter by extensions e.g. [\".docx\",\".pdf\"]."},
            },
        },
    },
    "search_documents": {
        "description": (
            "Search work documents by filename and/or content keywords in the local "
            "OneDrive for Business sync folder. No authentication required."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "query":          {"type": "string",  "description": "Keywords to search (all must match)."},
                "limit":          {"type": "integer", "description": "Max results (default 15, max 50).", "default": 15},
                "content_search": {"type": "boolean", "description": "Also search inside text/md files (default true).", "default": True},
                "extensions":     {"type": "array",   "items": {"type": "string"}, "description": "Filter by extensions."},
            },
            "required": ["query"],
        },
    },
    "read_document": {
        "description": (
            "Extract plain text from a local work document. "
            "Supports .txt, .md, .csv, .docx, .pptx, .xlsx, .pdf. "
            "Use 'path' from search_documents or list_recent_documents."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "path":      {"type": "string",  "description": "Absolute local path to the file."},
                "max_chars": {"type": "integer", "description": "Truncate at this many chars (default 8000).", "default": 8000},
            },
            "required": ["path"],
        },
    },
}

_TOOL_HANDLERS = {
    "list_recent_documents": _tool_list_recent_documents,
    "search_documents":      _tool_search_documents,
    "read_document":         _tool_read_document,
}


def _send(obj: dict) -> None:
    sys.stdout.write(json.dumps(obj) + "\n")
    sys.stdout.flush()


def _handle(request: dict) -> dict | None:
    method = request.get("method", "")
    req_id = request.get("id")

    if method == "initialize":
        return {
            "jsonrpc": "2.0", "id": req_id,
            "result": {
                "protocolVersion": "2024-11-05",
                "capabilities": {"tools": {}},
                "serverInfo": {"name": "artha-work-docs", "version": "2.0.0"},
            },
        }
    if method == "notifications/initialized":
        return None
    if method == "tools/list":
        return {"jsonrpc": "2.0", "id": req_id, "result": {"tools": [{"name": n, **d} for n, d in _TOOLS.items()]}}
    if method == "tools/call":
        params    = request.get("params", {})
        tool_name = params.get("name", "")
        tool_args = params.get("arguments", {})
        handler   = _TOOL_HANDLERS.get(tool_name)
        if not handler:
            return {"jsonrpc": "2.0", "id": req_id, "error": {"code": -32601, "message": f"Unknown tool: {tool_name}"}}
        try:
            result = handler(tool_args)
            return {
                "jsonrpc": "2.0", "id": req_id,
                "result": {"content": [{"type": "text", "text": json.dumps(result, indent=2)}], "isError": "error" in result},
            }
        except Exception as exc:
            _log.exception("Tool %s failed", tool_name)
            return {"jsonrpc": "2.0", "id": req_id, "result": {"content": [{"type": "text", "text": str(exc)}], "isError": True}}
    if method == "ping":
        return {"jsonrpc": "2.0", "id": req_id, "result": {}}
    if req_id is not None:
        return {"jsonrpc": "2.0", "id": req_id, "error": {"code": -32601, "message": f"Method not found: {method}"}}
    return None


def main() -> None:
    for line in sys.stdin:
        line = line.strip()
        if not line:
            continue
        try:
            request = json.loads(line)
        except json.JSONDecodeError as exc:
            _send({"jsonrpc": "2.0", "id": None, "error": {"code": -32700, "message": str(exc)}})
            continue
        response = _handle(request)
        if response is not None:
            _send(response)


if __name__ == "__main__":
    main()
